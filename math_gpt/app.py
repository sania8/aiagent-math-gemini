from flask import Flask, request, render_template, session, jsonify
import google.generativeai as genai
import plotly.graph_objects as go
import sympy as sp
import json
import re
import PIL.Image
import io
from PyPDF2 import PdfReader
from docx import Document
import pptx
app = Flask(__name__)
app.secret_key = "some_secret_key_for_session"


@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        api_key = request.form.get('api_key')
        if api_key:
            session['api_key'] = api_key
            genai.configure(api_key=api_key)
            return render_template('math_solver.html')
    return render_template('index.html')


@app.route('/extract_text', methods=['POST'])
def extract_text():
    if 'api_key' not in session:
        return jsonify({'error': 'API key missing'})

    genai.configure(api_key=session['api_key'])

    file = request.files.get('image')
    if not file:
        return jsonify({'error': 'No image provided'})

    image = PIL.Image.open(io.BytesIO(file.read()))

    return run_gemini_extraction(image)


@app.route('/solve', methods=['POST'])
def solve():
    if 'api_key' not in session:
        return jsonify({'error': 'API key missing'})

    genai.configure(api_key=session['api_key'])

    question = request.form.get('question')
    action = request.form.get('action')

    if not question:
        return jsonify({'error': 'Question is required'})

    model = genai.GenerativeModel("gemini-1.5-flash")

    if action == 'solve':
        prompt = f"Solve this maths question: {question}"
        response = model.generate_content(prompt)
        return jsonify({'result': response.text.strip()})

    elif action == 'explain':
        prompt = f"Explain step-by-step solution for: {question}"
        response = model.generate_content(prompt)
        return jsonify({'result': response.text.strip()})

    elif action == 'plot':
        return handle_plot_request(question, model)

    else:
        return jsonify({'error': 'Invalid action'})


@app.route('/upload_and_extract', methods=['POST'])
def upload_and_extract():
    if 'api_key' not in session:
        return jsonify({'error': 'API key missing'})

    genai.configure(api_key=session['api_key'])

    file = request.files.get('file')
    if not file:
        return jsonify({'error': 'No file provided'})

    filename = file.filename.lower()

    if filename.endswith(('jpg', 'jpeg', 'png')):
        return extract_from_image(file)
    elif filename.endswith('pdf'):
        return extract_from_pdf(file)
    elif filename.endswith('docx'):
        return extract_from_word(file)
    elif filename.endswith('pptx'):
        return extract_from_ppt(file)
    else:
        return jsonify({'error': 'Unsupported file type. Supported: JPG, PNG, PDF, DOCX, PPTX'})


def extract_from_image(file):
    image = PIL.Image.open(io.BytesIO(file.read()))
    return run_gemini_extraction(image)


def extract_from_pdf(file):
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return extract_math_from_text(text)


def extract_from_word(file):
    doc = Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    return extract_math_from_text(text)


def extract_from_ppt(file):
    presentation = pptx.Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return extract_math_from_text(text)


def extract_math_from_text(text):
    model = genai.GenerativeModel('gemini-1.5-flash')
    prompt = f"Extract only math questions (no explanations) from this text:\n{text}"
    response = model.generate_content(prompt)
    return jsonify({'text': response.text.strip()})


def run_gemini_extraction(image):
    model = genai.GenerativeModel('gemini-1.5-flash')
    response = model.generate_content([
        "Extract and return only math questions from this image.",
        image
    ])
    return jsonify({'text': response.text.strip()})


def extract_json_from_response(response_text):
    json_pattern = r'\{.*?\}'
    match = re.search(json_pattern, response_text, re.DOTALL)
    if not match:
        raise ValueError(f"Failed to locate JSON block in response:\n{response_text}")
    return json.loads(match.group(0))


def handle_plot_request(question, model):
    plot_prompt = f"""
    You are a math expert helping a student.
    The student asked this question: "{question}".

    Task:
    1. Extract the clean mathematical expression (pure math) from the question.
    2. Identify the type of equation. Choose from: 'explicit', 'parametric', 'implicit'.
    3. Decide if plotting a graph would help answer this question (true/false).

    Respond strictly in this JSON format:
    {{
        "expression": "extracted_math_expression",
        "type": "explicit/parametric/implicit",
        "graph_needed": true/false
    }}
    """

    response = model.generate_content(plot_prompt)
    response_text = response.text.strip()

    try:
        analysis = extract_json_from_response(response_text)

        if not analysis.get('graph_needed', False):
            return jsonify({'result': 'This question does not require a graph.'})

        expression = analysis.get('expression', '')
        expression_type = analysis.get('type', '')

        if not expression or expression_type not in {'explicit', 'parametric', 'implicit'}:
            return jsonify({'error': 'Invalid response from Gemini.'})

        graph_html = generate_plot_html(expression, expression_type)
        return render_template('plot.html', expression=expression, graph_html=graph_html)

    except Exception as e:
        return jsonify({'error': f'Failed to process Gemini response: {str(e)}\n\nResponse received:\n{response_text}'})


def generate_plot_html(expression, expression_type):
    expression = expression.replace('^', '**')
    x, y = sp.symbols('x y')

    fig = go.Figure()

    if expression_type == 'explicit':
        expr = sp.sympify(expression)
        x_vals = [i / 10 for i in range(-100, 101)]
        y_vals = [float(expr.subs(x, val)) for val in x_vals]
        fig.add_trace(go.Scatter(x=x_vals, y=y_vals, mode='lines'))

    # You can expand for parametric/implicit if needed.

    return fig.to_html(full_html=False)


if __name__ == '__main__':
    app.run(debug=True)
